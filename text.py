import os
import json
# from groq import Groq
from pypdf import PdfReader
from pptx import Presentation
from openai import OpenAI

# use openai client functions

# find out how to do file attachments

client = OpenAI(
    organization="org-xKM8WFGSFHsIgXs0hY7HsOKu",
    api_key=os.environ.get("OPENAI_API_KEY"),
)

with open("./dictionary.txt", 'r') as f:
    BRAINROT = f.read()

def extract_pdf_text(file_path):
    try:
        with open(file_path, 'rb') as file:
            reader = PdfReader(file)
            text = ''
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text
    except Exception as e:
        print(f"Error reading PDF: {e}")
        return ""

def extract_pptx_text(file_path):
    text_content = []
    try:
        # Open the presentation
        presentation = Presentation(file_path)
        
        # Loop through slides
        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_text = []
            # Loop through shapes in the slide
            for shape in slide.shapes:
                # Check if the shape has text
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        slide_text.append(paragraph.text.strip())
                # If the shape is a table, extract text from cells
                elif shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            slide_text.append(cell.text.strip())
                """
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
                """
            
            # Add slide number for clarity
            if slide_text:
                text_content.append(f"Slide {slide_number}:\n" + "\n".join(slide_text))
        
        return "\n\n".join(text_content)
    
    except Exception as e:
        print(f"Error extracting text: {e}")
        return ""

def extract_text(file_path):
    """
    Determine file type and extract text using the appropriate method.
    :param file_path: Path to the file (PDF or PPTX).
    :return: Extracted text as a string.
    """
    file_extension = os.path.splitext(file_path)[-1].lower()

    if file_extension == '.pdf':
        return extract_pdf_text(file_path)
    elif file_extension == '.pptx':
        return extract_pptx_text(file_path)
    else:
        raise Exception("Unsupported file type. Please provide a PDF or PPTX file.")


def generate_transcript(fp, skibidi_mode=True):
    """
    Generate a short-form video narration from a specified document.
    :param fp: Path to the file (PDF or PPTX).
    :param skibidi_mode: Use "brainrot" terms
    :return: Generate text as a string.
    """
    data = { "narration": "", "key_terms": [] }
    try:
        text_content = extract_text(fp)
        if text_content == "":
            return data
    except Exception as e:
        print(e)
        return data

    if skibidi_mode:
        system_prompt = f'"Brainrot" internet terms often refer to phrases, memes, or slang adopted by Generation Z and used excessively in online communities, sometimes to the point of being humorously overused or mindlessly repeated. Here is a dictionary of most of these terms:\n{BRAINROT}\n\n'
        system_prompt += """
            You are a short-form video creator FREQUENTLY using the terms FROM THE PROVIDED DICTIONARY when presenting information and explaining concepts, this also includes pop culture references, with a VERY STRICT EMPHASIS on MINDLESSLY INSERTING THROUGHTOUT YOUR DIALOGUE these "brainrot" words and slang: skibidi, rizz, Ohio, gyatt, sigma, based, bussy, gooning, etc. Your task will be to create a narration (no need for a greeting) summarizing and EXPLAINING IN BRAINROT TERMS the provided information from a document (most likely an educational material) that will be used for a TikTok or Instagram reel with a duration of 30 seconds to 1 minute AND make a list of ALL key terms (verbatim) mentioned in the narration. Responses WILL ONLY be in JSON using this schema
            {
                "narration": text string of the narration,
                "key_terms": [list of key terms]
            }
        """
    else:
        system_prompt = """
            You are a Gen Z short-form video creator. You will be writing a narration (no need for a greeting) summarizing and explaning the provided information from a document (most likely an educational material) that will be spoken by an AI voice for a faceless TikTok or Instagram reel with a duration of 30 seconds to 1 minute AND make a list of ALL key terms (verbatim) mentioned in the narration. Responses WILL ONLY be in JSON using this schema
            {
                "narration": text string of the narration,
                "key_terms": [list of key terms]
            }
        """

    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "system",
                    "content": system_prompt,
                },
                {
                    "role": "user",
                    "content": extract_pdf_text(fp),
                }
            ],
            response_format = { "type": "json_object" },
            # temperature=0.5
        )
        # print(response)
        return json.loads(response.choices[0].message.content)
    
    except Exception as e:
        print(f"Error communicating with OpenAI API: {e}")
        return data

if __name__ == "__main__":
    test_file = "C:/Users/abyam/Desktop/Research-Project/examples/mil-finals.pdf"
    print(generate_transcript(test_file, True))
